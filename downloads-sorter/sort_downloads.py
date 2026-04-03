import os
import sys
import shutil
import textwrap
from pathlib import Path

try:
    import anthropic
except ImportError:
    print("Installing anthropic SDK...")
    os.system("pip install anthropic")
    import anthropic

try:
    import pdfplumber
except ImportError:
    os.system("pip install pdfplumber")
    import pdfplumber

try:
    from docx import Document
except ImportError:
    os.system("pip install python-docx")
    from docx import Document

try:
    from pptx import Presentation
except ImportError:
    os.system("pip install python-pptx")
    from pptx import Presentation

try:
    import openpyxl
except ImportError:
    os.system("pip install openpyxl")
    import openpyxl


DOWNLOADS_DIR = Path.home() / "Downloads"
TEMP_DIR = DOWNLOADS_DIR / "temp"

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt",
    ".xlsx", ".xls", ".txt", ".md", ".csv",
    ".json", ".xml", ".yaml", ".yml", ".log",
    ".html", ".htm", ".rtf"
}

MAX_TEXT_LENGTH = 15000

SYSTEM_PROMPT = """You are a file summarization assistant. Given the content of a document, provide a concise summary in Chinese (2-4 sentences) that covers:
1. What type of document this is
2. The main topic or purpose
3. Key points or findings

Keep it brief and practical. Only output the summary, nothing else."""


def extract_text_pdf(filepath: Path) -> str:
    text_parts = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages[:10]:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
    except Exception as e:
        return f"[PDF extraction error: {e}]"
    return "\n".join(text_parts)


def extract_text_docx(filepath: Path) -> str:
    try:
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_text_pptx(filepath: Path) -> str:
    try:
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_text_xlsx(filepath: Path) -> str:
    try:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    texts.append(" | ".join(cells))
        wb.close()
        return "\n".join(texts[:500])
    except Exception as e:
        return f"[XLSX extraction error: {e}]"


def extract_text_plain(filepath: Path) -> str:
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"[Text extraction error: {e}]"


def extract_text(filepath: Path) -> str:
    ext = filepath.suffix.lower()
    extractors = {
        ".pdf": extract_text_pdf,
        ".docx": extract_text_docx,
        ".doc": extract_text_docx,
        ".pptx": extract_text_pptx,
        ".ppt": extract_text_pptx,
        ".xlsx": extract_text_xlsx,
        ".xls": extract_text_xlsx,
    }
    extractor = extractors.get(ext, extract_text_plain)
    text = extractor(filepath)
    if len(text) > MAX_TEXT_LENGTH:
        text = text[:MAX_TEXT_LENGTH] + "\n...[truncated]"
    return text


def summarize_file(client: anthropic.Anthropic, filepath: Path, content: str) -> str:
    try:
        message = client.messages.create(
            model="MiniMax-M2.7-highspeed",
            max_tokens=300,
            system=SYSTEM_PROMPT,
            messages=[{
                "role": "user",
                "content": f"File: {filepath.name}\n\nContent:\n{content}"
            }]
        )
        summary = ""
        for block in message.content:
            if block.type == "text":
                summary += block.text
        return summary.strip()
    except Exception as e:
        return f"[Summarization error: {e}]"


def main():
    api_key = os.environ.get("MINIMAX_API_KEY")
    if not api_key:
        api_key = input("Enter your MiniMax API key: ").strip()
        if not api_key:
            print("API key required. Set MINIMAX_API_KEY env var or enter it now.")
            sys.exit(1)

    os.environ["ANTHROPIC_BASE_URL"] = "https://api.minimaxi.com/anthropic"
    os.environ["ANTHROPIC_API_KEY"] = api_key
    client = anthropic.Anthropic()

    if not DOWNLOADS_DIR.exists():
        print(f"Downloads folder not found: {DOWNLOADS_DIR}")
        sys.exit(1)

    files = [
        f for f in DOWNLOADS_DIR.iterdir()
        if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS
    ]

    if not files:
        print("No supported files found in Downloads folder.")
        sys.exit(0)

    print(f"Found {len(files)} files to review.\n")

    kept = 0
    moved = 0
    errors = 0
    stopped = False

    for i, filepath in enumerate(files, 1):
        if stopped:
            break

        print(f"[{i}/{len(files)}] {filepath.name}")
        print(f"  Size: {filepath.stat().st_size / 1024:.1f} KB")

        content = extract_text(filepath)
        if not content or content.startswith("["):
            print(f"  {content}")
            errors += 1
            print()
            continue

        print("  Summarizing...")
        summary = summarize_file(client, filepath, content)
        print(f"  Summary: {summary}")

        while True:
            choice = input("  Keep? (y/n/skip/quit): ").strip().lower()
            if choice in ("y", "yes", ""):
                kept += 1
                print("  -> Kept in place")
                break
            elif choice in ("n", "no"):
                TEMP_DIR.mkdir(exist_ok=True)
                dest = TEMP_DIR / filepath.name
                counter = 1
                while dest.exists():
                    stem = filepath.stem
                    dest = TEMP_DIR / f"{stem}_{counter}{filepath.suffix}"
                    counter += 1
                shutil.move(str(filepath), str(dest))
                moved += 1
                print(f"  -> Moved to temp/")
                break
            elif choice in ("skip", "s"):
                print("  -> Skipped")
                break
            elif choice in ("quit", "q"):
                print("\nStopped by user.")
                stopped = True
                break
            else:
                print("  Invalid input. Use y/n/skip/quit.")

        print()

    print("\n" + "=" * 50)
    print("Report:")
    print(f"  Kept: {kept}")
    print(f"  Moved to temp/: {moved}")
    print(f"  Errors: {errors}")
    print(f"  Remaining in Downloads: {len([f for f in DOWNLOADS_DIR.iterdir() if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS])}")
    print("=" * 50)


if __name__ == "__main__":
    main()
