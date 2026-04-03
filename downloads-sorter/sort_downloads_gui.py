import os
import sys
import shutil
import threading
from pathlib import Path
from datetime import datetime

try:
    import customtkinter as ctk
except ImportError:
    os.system("pip install customtkinter")
    import customtkinter as ctk

try:
    import anthropic
except ImportError:
    os.system("pip install anthropic")
    import anthropic

try:
    import pdfplumber
except ImportError:
    os.system("pip install pdfplumber")
    import pdfplumber

try:
    from docx import Document
except ImportError:
    os.system("pip install python-docx")
    from docx import Document

try:
    from pptx import Presentation
except ImportError:
    os.system("pip install python-pptx")
    from pptx import Presentation

try:
    import openpyxl
except ImportError:
    os.system("pip install openpyxl")
    import openpyxl


DOWNLOADS_DIR = Path.home() / "Downloads"
TEMP_DIR = DOWNLOADS_DIR / "temp"
DEFAULT_KEEP_DIR = Path.home() / "Documents" / "Sorted"

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt",
    ".xlsx", ".xls", ".txt", ".md", ".csv",
    ".json", ".xml", ".yaml", ".yml", ".log",
    ".html", ".htm", ".rtf"
}

MAX_TEXT_LENGTH = 15000

SYSTEM_PROMPT = """You are a file summarization assistant. Given the content of a document, provide a concise summary in Chinese (2-4 sentences) that covers:
1. What type of document this is
2. The main topic or purpose
3. Key points or findings

Keep it brief and practical. Only output the summary, nothing else."""


def extract_text_pdf(filepath: Path) -> str:
    text_parts = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages[:10]:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
    except Exception as e:
        return f"[PDF extraction error: {e}]"
    return "\n".join(text_parts)


def extract_text_docx(filepath: Path) -> str:
    try:
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_text_pptx(filepath: Path) -> str:
    try:
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_text_xlsx(filepath: Path) -> str:
    try:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    texts.append(" | ".join(cells))
        wb.close()
        return "\n".join(texts[:500])
    except Exception as e:
        return f"[XLSX extraction error: {e}]"


def extract_text_plain(filepath: Path) -> str:
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"[Text extraction error: {e}]"


def extract_text(filepath: Path) -> str:
    ext = filepath.suffix.lower()
    extractors = {
        ".pdf": extract_text_pdf,
        ".docx": extract_text_docx,
        ".doc": extract_text_docx,
        ".pptx": extract_text_pptx,
        ".ppt": extract_text_pptx,
        ".xlsx": extract_text_xlsx,
        ".xls": extract_text_xlsx,
    }
    extractor = extractors.get(ext, extract_text_plain)
    text = extractor(filepath)
    if len(text) > MAX_TEXT_LENGTH:
        text = text[:MAX_TEXT_LENGTH] + "\n...[truncated]"
    return text


def summarize_file(client: anthropic.Anthropic, filepath: Path, content: str) -> str:
    try:
        message = client.messages.create(
            model="MiniMax-M2.7-highspeed",
            max_tokens=300,
            system=SYSTEM_PROMPT,
            messages=[{
                "role": "user",
                "content": f"File: {filepath.name}\n\nContent:\n{content}"
            }]
        )
        summary = ""
        for block in message.content:
            if block.type == "text":
                summary += block.text
        return summary.strip()
    except Exception as e:
        return f"[Summarization error: {e}]"


class SorterApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Downloads 整理助手")
        self.root.state("zoomed")

        self.api_key = os.environ.get("MINIMAX_API_KEY", "")
        self.client = None
        self.files = []
        self.current_index = 0
        self.current_file = None
        self.running = False
        self.api_visible = False

        self.stats = {"kept": 0, "moved": 0, "errors": 0, "skipped": 0}

        self._build_ui()

    def _build_ui(self):
        main = ctk.CTkFrame(self.root, fg_color="transparent")
        main.pack(fill=ctk.BOTH, expand=True, padx=28, pady=24)

        self._build_top_bar(main)

        ctk.CTkFrame(main, fg_color="#E5E5EA", height=2).pack(fill=ctk.X, pady=(14, 14))

        self._build_middle_section(main)

        ctk.CTkFrame(main, fg_color="#E5E5EA", height=2).pack(fill=ctk.X, pady=(14, 14))

        bottom_pane = ctk.CTkFrame(main, fg_color="transparent")
        bottom_pane.pack(fill=ctk.BOTH, expand=True)

        screen_w = self.root.winfo_screenwidth()
        half_width = int((screen_w - 56 - 14) / 2)

        left_pane = ctk.CTkFrame(bottom_pane, fg_color="transparent", width=half_width)
        left_pane.pack(side=ctk.LEFT, fill=ctk.BOTH, expand=True, padx=(0, 14))
        left_pane.pack_propagate(False)

        right_pane = ctk.CTkFrame(bottom_pane, fg_color="transparent")
        right_pane.pack(side=ctk.RIGHT, fill=ctk.BOTH, expand=True)

        self._build_left_pane(left_pane)
        self._build_right_pane(right_pane)

    def _build_top_bar(self, parent):
        frame = ctk.CTkFrame(parent, fg_color="transparent")
        frame.pack(fill=ctk.X)

        api_row = ctk.CTkFrame(frame, fg_color="transparent")
        api_row.pack(fill=ctk.X, pady=(0, 8))

        ctk.CTkLabel(api_row, text="API Key", font=("Segoe UI", 18, "bold"), text_color="#1D1D1F", width=100).pack(side=ctk.LEFT, padx=(0, 10))
        self.api_var = ctk.StringVar(value=self.api_key)
        self.api_entry = ctk.CTkEntry(api_row, textvariable=self.api_var, show="*", font=("Segoe UI", 17), corner_radius=10, height=40)
        self.api_entry.pack(side=ctk.LEFT, fill=ctk.X, expand=True, padx=(0, 10))
        ctk.CTkButton(api_row, text="显示", command=self._toggle_api, font=("Segoe UI", 15), corner_radius=10, width=70, height=40, fg_color="#FFFFFF", text_color="#0071E3", hover_color="#F0F0F2").pack(side=ctk.LEFT)

        scan_row = ctk.CTkFrame(frame, fg_color="transparent")
        scan_row.pack(fill=ctk.X, pady=(0, 8))

        ctk.CTkLabel(scan_row, text="扫描文件夹", font=("Segoe UI", 18, "bold"), text_color="#1D1D1F", width=100).pack(side=ctk.LEFT, padx=(0, 10))
        self.scan_folder_var = ctk.StringVar(value=str(DOWNLOADS_DIR))
        ctk.CTkEntry(scan_row, textvariable=self.scan_folder_var, font=("Segoe UI", 17), corner_radius=10, height=40).pack(side=ctk.LEFT, fill=ctk.X, expand=True, padx=(0, 10))
        ctk.CTkButton(scan_row, text="浏览", command=self._browse_scan_folder, font=("Segoe UI", 15), corner_radius=10, width=70, height=40, fg_color="#FFFFFF", text_color="#0071E3", hover_color="#F0F0F2").pack(side=ctk.LEFT)

        folder_row = ctk.CTkFrame(frame, fg_color="transparent")
        folder_row.pack(fill=ctk.X)

        left_folders = ctk.CTkFrame(folder_row, fg_color="transparent")
        left_folders.pack(side=ctk.LEFT, fill=ctk.X, expand=True, padx=(0, 10))

        ctk.CTkLabel(left_folders, text="保留文件夹", font=("Segoe UI", 18, "bold"), text_color="#1D1D1F", width=100).pack(side=ctk.LEFT, padx=(0, 10))
        self.keep_folder_var = ctk.StringVar(value=str(Path.home() / "Documents" / "Sorted"))
        ctk.CTkEntry(left_folders, textvariable=self.keep_folder_var, font=("Segoe UI", 17), corner_radius=10, height=40).pack(side=ctk.LEFT, fill=ctk.X, expand=True, padx=(0, 10))
        ctk.CTkButton(left_folders, text="浏览", command=self._browse_keep_folder, font=("Segoe UI", 15), corner_radius=10, width=70, height=40, fg_color="#FFFFFF", text_color="#0071E3", hover_color="#F0F0F2").pack(side=ctk.LEFT)

        right_folders = ctk.CTkFrame(folder_row, fg_color="transparent")
        right_folders.pack(side=ctk.RIGHT, fill=ctk.X, expand=True)

        ctk.CTkLabel(right_folders, text="拟删文件夹", font=("Segoe UI", 18, "bold"), text_color="#1D1D1F", width=100).pack(side=ctk.LEFT, padx=(0, 10))
        self.move_folder_var = ctk.StringVar(value=str(TEMP_DIR))
        ctk.CTkEntry(right_folders, textvariable=self.move_folder_var, font=("Segoe UI", 17), corner_radius=10, height=40).pack(side=ctk.LEFT, fill=ctk.X, expand=True, padx=(0, 10))
        ctk.CTkButton(right_folders, text="浏览", command=self._browse_move_folder, font=("Segoe UI", 15), corner_radius=10, width=70, height=40, fg_color="#FFFFFF", text_color="#0071E3", hover_color="#F0F0F2").pack(side=ctk.LEFT)

    def _toggle_api(self):
        self.api_visible = not self.api_visible
        self.api_entry.configure(show="" if self.api_visible else "*")

    def _browse_scan_folder(self):
        from tkinter import filedialog
        folder = filedialog.askdirectory(initialdir=self.scan_folder_var.get())
        if folder:
            self.scan_folder_var.set(folder)

    def _browse_keep_folder(self):
        from tkinter import filedialog
        folder = filedialog.askdirectory(initialdir=self.keep_folder_var.get())
        if folder:
            self.keep_folder_var.set(folder)

    def _browse_move_folder(self):
        from tkinter import filedialog
        folder = filedialog.askdirectory(initialdir=self.move_folder_var.get())
        if folder:
            self.move_folder_var.set(folder)

    def _build_middle_section(self, parent):
        card = ctk.CTkFrame(parent, fg_color="#FFFFFF", corner_radius=14)
        card.pack(fill=ctk.X)
        inner = ctk.CTkFrame(card, fg_color="transparent")
        inner.pack(fill=ctk.X, padx=20, pady=10)

        file_info_frame = ctk.CTkFrame(inner, fg_color="transparent")
        file_info_frame.pack(fill=ctk.X)

        self.file_name_var = ctk.StringVar(value="等待扫描...")
        self.file_meta_var = ctk.StringVar(value="")
        ctk.CTkLabel(file_info_frame, textvariable=self.file_name_var, font=("Segoe UI", 18, "bold"), text_color="#1D1D1F", anchor="w").pack(fill=ctk.X)
        ctk.CTkLabel(file_info_frame, textvariable=self.file_meta_var, font=("Segoe UI", 15), text_color="#86868B", anchor="w").pack(fill=ctk.X, pady=(4, 0))

    def _build_left_pane(self, parent):
        card = ctk.CTkFrame(parent, fg_color="#FFFFFF", corner_radius=14)
        card.pack(fill=ctk.BOTH, expand=True)
        inner = ctk.CTkFrame(card, fg_color="transparent")
        inner.pack(fill=ctk.BOTH, expand=True, padx=20, pady=18)

        ctk.CTkLabel(inner, text="AI 摘要", font=("Segoe UI", 18, "bold"), text_color="#1D1D1F", anchor="w").pack(fill=ctk.X, pady=(0, 10))
        self.summary_text = ctk.CTkTextbox(inner, font=("Segoe UI", 16), text_color="#1D1D1F", fg_color="#FAFAFA", corner_radius=10, border_width=1, border_color="#E5E5EA")
        self.summary_text.pack(fill=ctk.BOTH, expand=True)
        self.summary_text.configure(state=ctk.DISABLED)

    def _build_right_pane(self, parent):
        right_split = ctk.CTkFrame(parent, fg_color="transparent")
        right_split.pack(fill=ctk.BOTH, expand=True)

        btn_area = ctk.CTkFrame(right_split, fg_color="#FFFFFF", corner_radius=14)
        btn_area.pack(fill=ctk.X, pady=(0, 12))
        btn_inner = ctk.CTkFrame(btn_area, fg_color="transparent")
        btn_inner.pack(fill=ctk.X, padx=20, pady=14)

        btn_layout = ctk.CTkFrame(btn_inner, fg_color="transparent")
        btn_layout.pack(fill=ctk.X)

        btn_layout.grid_columnconfigure(0, weight=1)
        btn_layout.grid_columnconfigure(1, weight=3)

        self.btn_start = ctk.CTkButton(btn_layout, text="开始扫描", command=self._start_scan, font=("Segoe UI", 16, "bold"), corner_radius=10, height=64, fg_color="#0071E3", hover_color="#0077ED")
        self.btn_start.grid(row=0, column=0, rowspan=2, sticky="nsew", padx=(0, 14))

        btn_group = ctk.CTkFrame(btn_layout, fg_color="transparent")
        btn_group.grid(row=0, column=1, rowspan=2, sticky="nsew")
        btn_group.grid_rowconfigure(0, weight=1)
        btn_group.grid_rowconfigure(1, weight=1)
        btn_group.grid_columnconfigure(0, weight=1)
        btn_group.grid_columnconfigure(1, weight=1)

        self.btn_keep = ctk.CTkButton(btn_group, text="保留", command=lambda: self._action("keep"), font=("Segoe UI", 16, "bold"), corner_radius=10, height=28, fg_color="#34C759", hover_color="#2DB84E", state=ctk.DISABLED)
        self.btn_keep.grid(row=0, column=0, sticky="nsew", padx=(0, 8), pady=(0, 4))
        self.btn_move = ctk.CTkButton(btn_group, text="拟删", command=lambda: self._action("move"), font=("Segoe UI", 16, "bold"), corner_radius=10, height=28, fg_color="#FF3B30", hover_color="#E0342B", state=ctk.DISABLED)
        self.btn_move.grid(row=0, column=1, sticky="nsew", padx=(8, 0), pady=(0, 4))

        self.btn_skip = ctk.CTkButton(btn_group, text="跳过", command=lambda: self._action("skip"), font=("Segoe UI", 16), corner_radius=10, height=28, fg_color="#F0F0F2", text_color="#1D1D1F", hover_color="#E5E5EA", state=ctk.DISABLED)
        self.btn_skip.grid(row=1, column=0, sticky="nsew", padx=(0, 8), pady=(4, 0))
        self.btn_quit = ctk.CTkButton(btn_group, text="停止", command=lambda: self._action("quit"), font=("Segoe UI", 16), corner_radius=10, height=28, fg_color="#F0F0F2", text_color="#86868B", hover_color="#E5E5EA", state=ctk.DISABLED)
        self.btn_quit.grid(row=1, column=1, sticky="nsew", padx=(8, 0), pady=(4, 0))

        log_card = ctk.CTkFrame(right_split, fg_color="#FFFFFF", corner_radius=14)
        log_card.pack(fill=ctk.BOTH, expand=True)
        log_inner = ctk.CTkFrame(log_card, fg_color="transparent")
        log_inner.pack(fill=ctk.BOTH, expand=True, padx=20, pady=10)

        ctk.CTkLabel(log_inner, text="日志", font=("Segoe UI", 16, "bold"), text_color="#1D1D1F", anchor="w").pack(fill=ctk.X, pady=(0, 6))
        self.log_text = ctk.CTkTextbox(log_inner, font=("Consolas", 15), text_color="#86868B", fg_color="#FAFAFA", corner_radius=10, border_width=1, border_color="#E5E5EA")
        self.log_text.pack(fill=ctk.BOTH, expand=True)
        self.log_text.configure(state=ctk.DISABLED)

    def _log(self, msg):
        self.log_text.configure(state=ctk.NORMAL)
        self.log_text.insert(ctk.END, msg + "\n")
        self.log_text.see(ctk.END)
        self.log_text.configure(state=ctk.DISABLED)

    def _update_summary(self, text):
        self.summary_text.configure(state=ctk.NORMAL)
        self.summary_text.delete(1.0, ctk.END)
        self.summary_text.insert(ctk.END, text)
        self.summary_text.configure(state=ctk.DISABLED)

    def _set_buttons(self, running=False, processing=False):
        state = ctk.NORMAL if processing else ctk.DISABLED
        for btn in [self.btn_keep, self.btn_move, self.btn_skip, self.btn_quit]:
            btn.configure(state=state)
        start_state = ctk.DISABLED if processing else ctk.NORMAL
        self.btn_start.configure(state=start_state, text="开始扫描" if not running else "继续扫描")

    def _start_scan(self):
        api_key = self.api_var.get().strip()
        if not api_key:
            from tkinter import messagebox
            messagebox.showerror("错误", "请输入 MiniMax API Key", parent=self.root)
            return

        self.api_key = api_key
        os.environ["ANTHROPIC_BASE_URL"] = "https://api.minimaxi.com/anthropic"
        os.environ["ANTHROPIC_API_KEY"] = api_key
        try:
            self.client = anthropic.Anthropic()
        except Exception as e:
            from tkinter import messagebox
            messagebox.showerror("错误", f"初始化客户端失败: {e}", parent=self.root)
            return

        scan_folder = Path(self.scan_folder_var.get())
        self.files = [f for f in scan_folder.iterdir() if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS]
        if not self.files:
            from tkinter import messagebox
            messagebox.showinfo("提示", f"{scan_folder} 文件夹中没有找到支持的文档文件", parent=self.root)
            return
        self.files.sort(key=lambda f: f.stat().st_mtime, reverse=True)
        self._log(f"找到 {len(self.files)} 个文件，按修改时间倒序排列")

        self.current_index = 0
        self.stats = {"kept": 0, "moved": 0, "errors": 0, "skipped": 0}
        self.running = True
        self._set_buttons(running=True, processing=True)
        threading.Thread(target=self._process_files, daemon=True).start()

    def _process_files(self):
        try:
            while self.current_index < len(self.files) and self.running:
                self.current_file = self.files[self.current_index]
                self.root.after(0, self._update_ui_for_file)

                content = extract_text(self.current_file)
                if not content or content.startswith("["):
                    self.root.after(0, lambda: self._log(f"[错误] {self.current_file.name}: {content}"))
                    self.stats["errors"] += 1
                    self.root.after(0, self._next_file)
                    continue

                self.root.after(0, lambda: self._log(f"正在总结: {self.current_file.name}"))
                summary = summarize_file(self.client, self.current_file, content)
                self.root.after(0, lambda s=summary: self._update_summary(s))
                self.root.after(0, lambda: self._log(f"摘要完成: {self.current_file.name}"))
                self.root.after(0, lambda: self._set_buttons(running=True, processing=True))
                self._wait_for_action()

        except Exception as e:
            self.root.after(0, lambda: self._log(f"[异常] {e}"))
        finally:
            self.running = False
            self.root.after(0, self._finish)

    def _update_ui_for_file(self):
        size_kb = self.current_file.stat().st_size / 1024
        mtime_str = datetime.fromtimestamp(self.current_file.stat().st_mtime).strftime("%Y-%m-%d %H:%M")
        self.file_name_var.set(self.current_file.name)
        self.file_meta_var.set(f"{size_kb:.1f} KB · 修改于 {mtime_str} · [{self.current_index + 1}/{len(self.files)}]")

    def _next_file(self):
        self.current_index += 1
        if self.current_index < len(self.files):
            self.root.after(0, self._update_ui_for_file)
        else:
            self.running = False

    def _wait_for_action(self):
        self.action_result = None
        while self.action_result is None and self.running:
            self.root.update()
            import time
            time.sleep(0.1)

    def _action(self, action):
        self.action_result = action
        self._set_buttons(running=True, processing=True)
        name = self.current_file.name

        if action == "keep":
            dest_folder = Path(self.keep_folder_var.get())
            dest_folder.mkdir(parents=True, exist_ok=True)
            dest = dest_folder / name
            counter = 1
            while dest.exists():
                dest = dest_folder / f"{self.current_file.stem}_{counter}{self.current_file.suffix}"
                counter += 1
            shutil.move(str(self.current_file), str(dest))
            self.stats["kept"] += 1
            self._log(f"-> 保留至: {dest}")
        elif action == "move":
            move_folder = Path(self.move_folder_var.get())
            move_folder.mkdir(parents=True, exist_ok=True)
            dest = move_folder / name
            counter = 1
            while dest.exists():
                dest = TEMP_DIR / f"{self.current_file.stem}_{counter}{self.current_file.suffix}"
                counter += 1
            shutil.move(str(self.current_file), str(dest))
            self.stats["moved"] += 1
            self._log(f"-> 移到拟删文件夹: {name}")
        elif action == "skip":
            self.stats["skipped"] += 1
            self._log(f"-> 跳过: {name}")
        elif action == "quit":
            self.running = False
            self._log("-> 用户退出")
            return

        self._next_file()

    def _finish(self):
        self._set_buttons(running=False, processing=False)
        report = f"\n===== 报告 =====\n保留: {self.stats['kept']}\n移到 temp/: {self.stats['moved']}\n跳过: {self.stats['skipped']}\n错误: {self.stats['errors']}\n================"
        self._log(report)
        self.file_name_var.set("处理完成")
        self.file_meta_var.set("")
        self._update_summary("处理完成。查看日志获取详细报告。")


def main():
    ctk.set_appearance_mode("light")
    ctk.set_default_color_theme("blue")
    root = ctk.CTk()
    try:
        from ctypes import windll
        windll.shcore.SetProcessDpiAwareness(1)
    except Exception:
        pass
    SorterApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
